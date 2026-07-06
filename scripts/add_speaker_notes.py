"""Inject speaker notes (from the talk transcript) into the Red Hat Summit 2026
deck — one block of notes per slide. Idempotent: re-running replaces notes.
"""
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
DST = ROOT / "Media Kit" / "slide" / "Red Hat summit" / "Lessons Learned Orchestrating Multi-Tenant GPUs on OpenShift AI with NVIDIA KAI H200 - Luca Berton Red Hat Summit 2026.pptx"

# 19 entries — one per slide, in order.
NOTES = [
    # 1 — Title
    (
        "Multi-tenant GPUs on bare metal OpenShift AI. Today I want to share a "
        "practical blueprint for running multi-tenant GPUs on bare metal OpenShift "
        "AI, delivering a sovereign full-stack AI platform that drives immediate "
        "results. The key point: getting GPUs to run is not the hard part — the "
        "hard part is sharing them safely, fairly, and efficiently across teams.\n\n"
        "My name is Luca Berton, and I work on AI-ready infrastructure and "
        "enterprise cloud environments. This talk comes from hands-on platform "
        "work delivered to an avionics customer in France, in production today. "
        "What I will show you is the result of operational lessons from real "
        "systems."
    ),
    # 2 — The Environment
    (
        "This was a bare metal environment: OpenShift, OpenShift AI, NVIDIA GPU "
        "and Network Operators. H200 GPU nodes, RDMA storage, and an air-gapped "
        "setup with local mirrors. In the cloud, some complexity is abstracted "
        "away. Here, every layer is ours to own, configure, and recover."
    ),
    # 3 — "It runs" ≠ "It's safe to share"
    (
        "We learned very quickly that a working GPU cluster is not the same as a "
        "shareable GPU platform. The real risks were noisy neighbors, queue "
        "starvation, mismatched MIG configuration, driver drift, and networking "
        "failures around SR-IOV and RDMA. So the challenge became: how do we "
        "prevent chaos before tenants even feel it?"
    ),
    # 4 — The Framework (Safe / Fair / Efficient)
    (
        "Every platform decision had to pass through three lenses. First, SAFE: "
        "one team cannot break another. Second, FAIR: contention is "
        "deterministic, not political. Third, EFFICIENT: we optimize for useful "
        "outcomes per GPU-hour, not just utilization graphs."
    ),
    # 5 — Three Personas, One Platform
    (
        "We designed for three personas. The end user wants fast access and "
        "stable latency. The LLMOps teams want repeatable deployments, safe "
        "upgrades, and observability. The tenant admin wants boundaries and "
        "visibility into cost. If any one of those personas loses, the platform "
        "adoption collapses."
    ),
    # 6 — GitOps: Everything is Code, Auditable
    (
        "Our answer was GitOps. Argo CD and Kustomize became the single source "
        "of truth for operators, infrastructure, tenant configuration, and "
        "applications. The reason this matters is not just automation — it is "
        "auditability, repeatability, and rollback. In a multi-tenant GPU "
        "platform, manual fixes become future outages."
    ),
    # 7 — Bootstrap: bare metal to GitOps
    (
        "We used Ansible for the initial handshake, especially in the air-gapped "
        "environment, and then handed control to GitOps as early as possible. "
        "The goal was simple: bootstrap once, then let Argo CD own day-two "
        "operations. That reduces drift and makes the platform behavior "
        "predictable over time."
    ),
    # 8 — Multi-tenant networking: HAProxy + Keepalived
    (
        "Multi-tenant networking sat on HAProxy with Keepalived providing VRRP "
        "for the virtual IPs, and SR-IOV virtual functions giving each tenant "
        "its own NIC isolation and per-tenant network QoS. All of this is "
        "templated with Jinja2 so adding a tenant is one dictionary entry — and "
        "per-tenant logging flows into rsyslog with one file per tenant for "
        "auditing and troubleshooting."
    ),
    # 9 — Safety: reduce blast radius by design
    (
        "For us, safety means hard isolation boundaries: namespace isolation, "
        "scoped service accounts, least-privilege RBAC, deny-by-default network "
        "policies, pod security guardrails, quotas, limits, and admission "
        "checks. The key idea is that safety is not something operators "
        "remember to do — it is something the platform enforces by default."
    ),
    # 10 — Safety: tenant bootstrap bundle
    (
        "Each tenant came from a standard Kustomize build deployed through "
        "Argo CD. That bundle included namespace configuration, RBAC, "
        "NetworkPolicy, quotas, and supporting networking pieces. The practical "
        "outcome is important: no tickets, no manual steps, no tribal "
        "knowledge. A single Git pull request becomes a tenant provisioning "
        "workflow."
    ),
    # 11 — SR-IOV: NVIDIA Network Op vs NICs
    (
        "One subtle lesson was that bare metal forces you to care about the "
        "actual hardware mix. Different NICs required different operational "
        "models. NVIDIA ConnectX paths were managed by the NVIDIA Network "
        "Operator for RDMA and GPUDirect-capable flows, while other NIC paths "
        "were handled by the OpenShift SR-IOV Network Operator for management "
        "traffic. The main takeaway: do not assume one networking recipe fits "
        "every GPU node."
    ),
    # 12 — Open kernel modules + DMA-BUF
    (
        "We reduced upgrade fragility by moving away from the more tightly "
        "coupled legacy approach toward open kernel modules and DMA-BUF. The "
        "important message here is not the terminology — it is that safer "
        "upgrades require reducing coupling between drivers, kernel "
        "dependencies, and data paths."
    ),
    # 13 — Fairness: make Contention Deterministic
    (
        "In shared GPU environments, fairness does not happen naturally. "
        "Without clear rules, the loudest team wins. So we used hard tenant GPU "
        "caps, priority classes, explicit preemption posture, scheduling "
        "constraints, and GPU-aware scheduling visibility through KAI. The goal "
        "was simple: make resource contention predictable, explainable, and "
        "intentional."
    ),
    # 14 — Efficiency: outcomes per GPU-hour
    (
        "We looked at efficiency as outcomes per GPU-hour. That means choosing "
        "between time slicing, MIG, and full GPU based on workload type rather "
        "than ideology. Interactive notebooks, training, and inference each "
        "behave differently. The platform has to optimize for real workload "
        "outcomes, not just maximum theoretical utilization."
    ),
    # 15 — The upgrade plan IS the platform
    (
        "One of the most important lessons was that the upgrade plan is not an "
        "appendix — it is part of the platform design. We managed a known-good "
        "compatibility matrix in Git, used canary pools, bake periods, "
        "validation gates, and relied on Git-based rollback through Argo CD. In "
        "practice, trust in the platform depends heavily on trust in the "
        "upgrade path."
    ),
    # 16 — Make the invisible visible
    (
        "We invested in per-tenant monitoring, chargeback views, scheduler "
        "reports, GPU-hour visibility, queue time, latency, and network "
        "health. The big change was behavioral: once teams can see their own "
        "usage and impact, they become much better platform citizens. "
        "Visibility turns platform governance into something collaborative "
        "instead of punitive."
    ),
    # 17 — Demo
    (
        "In this demo, I will show you the usage pattern of three different "
        "personas.\n\n"
        "WILLIAM — LLMOps engineer. His goal is to deploy LLM models "
        "on-premises so developers can securely integrate their features into "
        "applications. William identifies a model on Hugging Face, chooses "
        "Microsoft Phi-4-mini-instruct, logs into OpenShift, and navigates to "
        "the OpenShift AI dashboard. After verifying the model files are stored "
        "in an object storage bucket, he configures the model deployment, sets "
        "up the model server using vLLM, and allocates the NVIDIA GPU for "
        "acceleration. He chooses between profile types based on the resources "
        "dedicated for this deployment. Once deployed, he retrieves the API "
        "endpoint and reviews dashboards: GPU accelerator and resources "
        "assigned, inference benchmarks (inter-token latency, throughput), "
        "comparisons with other models, and the OpenShift system metrics for "
        "CPU, memory, and resource traces.\n\n"
        "MICHAEL — OpenShift administrator. His goal is to guarantee "
        "performance, availability, and optimal utilization of the platform's "
        "GPU resources. Michael checks the OpenShift AI dashboard to see the "
        "active model deployments — including William's, and others such as a "
        "GPT-2. He then switches to Grafana to view the NVIDIA DCGM dashboard, "
        "monitoring power usage, temperature, SM clock, memory clock, and "
        "specific memory usage per MIG profile. Seven out of eight GPUs are in "
        "use.\n\n"
        "MARK — Python developer. His goal is to access an internal AI-powered "
        "coding assistant to boost productivity while remaining compliant with "
        "internal security policies. Mark works inside Visual Studio Code. He "
        "installs the open-source Continue AI code assistant extension and "
        "edits its config.yaml to point at the internal API endpoint William "
        "deployed earlier — not a public cloud service. He then opens a Python "
        "file with a sorting algorithm and asks the local LLM via Continue to "
        "'make this more readable', successfully generating refined code with "
        "diff lines in green and red."
    ),
    # 18 — Seven Guardrails You Can Apply Today
    (
        "If you take only a few things away: start with GitOps first and a "
        "standard tenant template. Those two alone remove a large amount of "
        "operational pain. Then add quotas, clear networking ownership, a "
        "safer GPU stack, useful observability, and an explicit upgrade "
        "playbook. The broader message: mature shared GPU platforms are built "
        "from guardrails, not heroics."
    ),
    # 19 — Closing
    (
        "You can get the full slide deck by scanning the QR code. Multi-tenant "
        "GPUs on bare metal work when safety is provable, fairness is "
        "explicit, and efficiency is measured. In our case, GitOps was the "
        "mechanism that made those principles enforceable.\n\n"
        "Thank you."
    ),
]


def set_notes(slide, text):
    """Replace the notes text frame content with `text`, preserving paragraph
    breaks (split on \n)."""
    notes = slide.notes_slide  # creates a notes slide if missing
    tf = notes.notes_text_frame
    tf.clear()
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # clear any pre-existing runs
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        run = p.add_run()
        run.text = line


def main():
    prs = Presentation(str(DST))
    n = len(prs.slides)
    if n != len(NOTES):
        print(f"WARN: deck has {n} slides, but {len(NOTES)} notes blocks provided")
    for i, slide in enumerate(prs.slides):
        if i >= len(NOTES):
            break
        set_notes(slide, NOTES[i])
        preview = NOTES[i].split("\n")[0][:70]
        print(f"  slide {i+1:2d}: {preview}…")
    prs.save(str(DST))
    print(f"\nSaved → {DST.name}")


if __name__ == "__main__":
    main()
