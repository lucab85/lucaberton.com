"""Build the social-media + recent-events Media Kit PPTX for Luca Berton.

Run: python3 scripts/build_media_kit_pptx.py
Output: Media Kit/Luca Berton Media Kit - Social & Events.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "Media Kit" / "Luca Berton Media Kit - Social & Events.pptx"
LOGO = ROOT / "Media Kit" / "LucaLogo" / "Luca_Berton_Full_Logo.png"
ICON = ROOT / "Media Kit" / "LucaLogo" / "Luca_Berton_Icon.png"

# Brand palette
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT = RGBColor(0x1F, 0x6F, 0xEB)   # blue
ACCENT2 = RGBColor(0xF5, 0xA6, 0x23)  # warm accent
LIGHT = RGBColor(0xF4, 0xF6, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x55, 0x5F, 0x6D)
DARK = RGBColor(0x10, 0x18, 0x28)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=DARK, bullet="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def header_bar(slide, title, kicker=None):
    add_rect(slide, 0, 0, SW, Inches(1.0), NAVY)
    add_rect(slide, 0, Inches(1.0), SW, Inches(0.06), ACCENT)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(11), Inches(0.6),
             title, size=28, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        add_text(slide, Inches(0.5), Inches(0.62), Inches(11), Inches(0.35),
                 kicker, size=12, color=RGBColor(0xC9, 0xD7, 0xF0))
    # small icon top-right
    if ICON.exists():
        slide.shapes.add_picture(str(ICON), SW - Inches(0.95), Inches(0.18),
                                 height=Inches(0.65))


def footer(slide, page_no=None, total=None):
    add_rect(slide, 0, SH - Inches(0.35), SW, Inches(0.35), LIGHT)
    add_text(slide, Inches(0.5), SH - Inches(0.33), Inches(8), Inches(0.3),
             "Luca Berton — Media Kit · Social & Recent Events · 2026",
             size=9, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    if page_no:
        add_text(slide, SW - Inches(2.5), SH - Inches(0.33), Inches(2),
                 Inches(0.3), f"{page_no} / {total}",
                 size=9, color=GREY, align=PP_ALIGN.RIGHT,
                 anchor=MSO_ANCHOR.MIDDLE)


def stat_card(slide, x, y, w, h, value, label, sub=None,
              fill=WHITE, accent=ACCENT):
    add_rect(slide, x, y, w, h, fill, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(slide, x, y, Inches(0.12), h, accent)
    add_text(slide, x + Inches(0.3), y + Inches(0.18), w - Inches(0.4),
             Inches(0.7), value, size=28, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.3), y + Inches(0.85), w - Inches(0.4),
             Inches(0.4), label, size=12, bold=True, color=DARK)
    if sub:
        add_text(slide, x + Inches(0.3), y + Inches(1.2), w - Inches(0.4),
                 h - Inches(1.25), sub, size=10, color=GREY)


# ---------------------------------------------------------------------------
# Slide 1 — Cover
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
# Decorative accent bars
add_rect(s, 0, Inches(6.7), SW, Inches(0.12), ACCENT)
add_rect(s, Inches(0.5), Inches(6.95), Inches(2.4), Inches(0.06), ACCENT2)

if LOGO.exists():
    s.shapes.add_picture(str(LOGO), Inches(0.55), Inches(0.5),
                         height=Inches(0.9))

add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(1.2),
         "Luca Berton", size=60, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(3.1), Inches(12), Inches(0.7),
         "Media Kit — Social Reach & Recent Event Coverage",
         size=24, color=RGBColor(0xC9, 0xD7, 0xF0))
add_text(s, Inches(0.6), Inches(3.85), Inches(12), Inches(0.5),
         "AI & Cloud Advisor · Docker Captain · KubeCon EU 2026 Speaker",
         size=16, color=ACCENT2)

add_text(
    s, Inches(0.6), Inches(4.6), Inches(12), Inches(1.3),
    "Featured at KubeCon Europe, FOSDEM and DevWorld — bringing\n"
    "platform engineering, GPU/MLOps and open-source insight to global stages.",
    size=15, color=WHITE,
)

add_text(s, Inches(0.6), Inches(5.95), Inches(12), Inches(0.4),
         "lucaberton.com  ·  linkedin.com/in/lucaberton  ·  youtube.com/@BertonLuca  ·  x.com/TheLucaBerton",
         size=12, color=RGBColor(0xC9, 0xD7, 0xF0))
add_text(s, Inches(0.6), Inches(7.05), Inches(12), Inches(0.35),
         "Amsterdam, NL  ·  luca@lucaberton.com  ·  +31 6 1225 5399",
         size=11, color=RGBColor(0x9F, 0xB1, 0xCF))

# ---------------------------------------------------------------------------
# Slide 2 — At a Glance
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header_bar(s, "At a Glance",
           "Cloud-native engineer, author and conference speaker")

add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.0),
         "Luca Berton helps enterprises ship secure, scalable AI & Kubernetes "
         "platforms — and shares the journey across stages, podcasts and a "
         "fast-growing technical audience.",
         size=15, color=DARK)

cards = [
    ("18+", "Years in tech", "Enterprise platforms, automation, MLOps"),
    ("8",   "Technical books", "Apress · Springer · BPB"),
    ("30+", "Conference talks", "KubeCon · FOSDEM · DevWorld · Red Hat"),
    ("150k+", "Monthly views", "ansiblepilot.com + YouTube"),
]
cw = Inches(2.95)
gap = Inches(0.2)
x0 = Inches(0.5)
y0 = Inches(2.6)
ch = Inches(1.85)
for i, (v, l, sub) in enumerate(cards):
    stat_card(s, x0 + i * (cw + gap), y0, cw, ch, v, l, sub)

add_rect(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.2), LIGHT)
add_text(s, Inches(0.75), Inches(4.85), Inches(11.8), Inches(0.4),
         "Signature Topics", size=14, bold=True, color=NAVY)
add_bullets(
    s, Inches(0.75), Inches(5.25), Inches(5.8), Inches(1.6),
    [
        "Multi-tenant GPU orchestration on Kubernetes / OpenShift AI",
        "Pragmatic MLOps: from POC to production",
        "Platform engineering at million-resource scale",
    ],
    size=12,
)
add_bullets(
    s, Inches(6.7), Inches(5.25), Inches(6.0), Inches(1.6),
    [
        "Ansible & Terraform automation for the enterprise",
        "Open source community building (Ansible Pilot)",
        "Secure, SOC2 / ISO 27001-aligned cloud architectures",
    ],
    size=12,
)
footer(s, 2, 10)

# ---------------------------------------------------------------------------
# Slide 3 — Social Media Reach
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header_bar(s, "Social Media Reach",
           "Engaged community across LinkedIn, YouTube, X, Instagram & Web")

socials = [
    ("LinkedIn",  "linkedin.com/in/lucaberton",
     "Daily insight on AI, Kubernetes, GPU & platform engineering",
     RGBColor(0x0A, 0x66, 0xC2)),
    ("YouTube · @BertonLuca", "youtube.com/@BertonLuca",
     "Hands-on Ansible, Kubernetes & cloud-native walkthroughs",
     RGBColor(0xFF, 0x00, 0x00)),
    ("X (Twitter)", "x.com/TheLucaBerton",
     "Live conference notes, hot takes and open-source signal",
     RGBColor(0x00, 0x00, 0x00)),
    ("Instagram", "instagram.com/thelucaberton",
     "Behind the scenes from talks, books and travel",
     RGBColor(0xE1, 0x30, 0x6C)),
    ("Ansible Pilot", "ansiblepilot.com  ·  ~150k views / month",
     "300+ Ansible & Terraform use cases, blog and YouTube hub",
     RGBColor(0xCC, 0x00, 0x00)),
    ("Website & Blog", "lucaberton.com",
     "Long-form articles, conference recaps, advisory & books",
     ACCENT),
]
x0 = Inches(0.5)
y0 = Inches(1.35)
cw = Inches(6.15)
ch = Inches(1.55)
gap = Inches(0.2)
for i, (name, handle, desc, color) in enumerate(socials):
    col = i % 2
    row = i // 2
    x = x0 + col * (cw + gap)
    y = y0 + row * (ch + Inches(0.18))
    add_rect(s, x, y, cw, ch, WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, x, y, Inches(0.16), ch, color)
    add_text(s, x + Inches(0.35), y + Inches(0.15), cw - Inches(0.5),
             Inches(0.45), name, size=16, bold=True, color=NAVY)
    add_text(s, x + Inches(0.35), y + Inches(0.6), cw - Inches(0.5),
             Inches(0.4), handle, size=12, bold=True, color=color)
    add_text(s, x + Inches(0.35), y + Inches(1.0), cw - Inches(0.5),
             Inches(0.55), desc, size=11, color=GREY)

footer(s, 3, 10)

# ---------------------------------------------------------------------------
# Slide 4 — Audience Snapshot
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header_bar(s, "Audience Snapshot",
           "Who is reading, watching and showing up at the talks")

add_text(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.5),
         "Audience Profile", size=18, bold=True, color=NAVY)
add_bullets(
    s, Inches(0.5), Inches(1.85), Inches(6.0), Inches(3.5),
    [
        "Platform & DevOps engineers, SREs, MLOps practitioners",
        "Cloud architects and engineering managers (AWS / Azure / GCP)",
        "CTOs and tech leads scaling AI / Kubernetes adoption",
        "Open-source contributors in the Ansible & CNCF ecosystems",
        "Technical authors, instructors and conference organisers",
    ],
    size=13,
)

add_text(s, Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.5),
         "Geography & Channels", size=18, bold=True, color=NAVY)
add_bullets(
    s, Inches(7.0), Inches(1.85), Inches(6.0), Inches(3.5),
    [
        "Strong presence in EMEA — NL, UK, DE, BE, IT, FR",
        "Growing reach in North America and India (book audiences)",
        "Long-form on LinkedIn & lucaberton.com",
        "Tutorial video on YouTube @BertonLuca",
        "Live event signal on X / @TheLucaBerton",
    ],
    size=13,
)

add_rect(s, Inches(0.5), Inches(5.45), Inches(12.3), Inches(1.55), LIGHT)
add_text(s, Inches(0.75), Inches(5.55), Inches(11.8), Inches(0.4),
         "Why brands collaborate", size=14, bold=True, color=NAVY)
add_text(
    s, Inches(0.75), Inches(5.95), Inches(11.8), Inches(1.05),
    "Independent practitioner voice · enterprise credibility (Dell, JPMorgan, Red Hat) · "
    "consistent multi-channel coverage of flagship cloud-native events · long-tail SEO via "
    "300+ how-to articles · authentic community building around Ansible Pilot.",
    size=12, color=DARK,
)
footer(s, 4, 10)

# ---------------------------------------------------------------------------
# Slide 5 — Recent Event Coverage Overview
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header_bar(s, "Recent Event Coverage",
           "Three flagship 2025–2026 stages — and how they were amplified")

events = [
    ("KubeCon + CloudNativeCon\nEurope 2026",
     "Amsterdam 🇳🇱  ·  March 2026",
     "Speaker — Lessons Learned Orchestrating Multi-Tenant\n"
     "GPUs on OpenShift AI with NVIDIA KAI (G/H200).\n"
     "MC at Cloud Native Rejekts EU.",
     ACCENT),
    ("FOSDEM 2026",
     "Brussels 🇧🇪  ·  Jan 31 – Feb 1, 2026",
     "Community attendee & connector — RISC-V, edge AI,\n"
     "open silicon. 8,000+ attendees, 1,013 talks.",
     ACCENT2),
    ("DevWorld Conference",
     "Amsterdam 🇳🇱  ·  Feb 2025",
     "Attendee & networker at one of Europe's largest\n"
     "developer festivals — AI, cloud, dev experience.",
     RGBColor(0x16, 0xA3, 0x4A)),
]
x0 = Inches(0.5)
y0 = Inches(1.35)
cw = Inches(4.15)
gap = Inches(0.15)
ch = Inches(4.2)
for i, (title, when, body, color) in enumerate(events):
    x = x0 + i * (cw + gap)
    add_rect(s, x, y0, cw, ch, WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, x, y0, cw, Inches(0.18), color)
    add_text(s, x + Inches(0.25), y0 + Inches(0.35), cw - Inches(0.4),
             Inches(1.0), title, size=18, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y0 + Inches(1.45), cw - Inches(0.4),
             Inches(0.4), when, size=11, bold=True, color=color)
    add_text(s, x + Inches(0.25), y0 + Inches(1.95), cw - Inches(0.4),
             Inches(2.2), body, size=12, color=DARK)

add_rect(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.3), LIGHT)
add_text(s, Inches(0.75), Inches(5.8), Inches(11.8), Inches(0.4),
         "Coverage formats produced", size=14, bold=True, color=NAVY)
add_text(
    s, Inches(0.75), Inches(6.2), Inches(11.8), Inches(0.8),
    "Long-form blog recap · LinkedIn carousel & posts · YouTube vlog & interviews · "
    "X live thread · podcast guest appearances · session slides on lucaberton.com/slides.",
    size=12, color=DARK,
)
footer(s, 5, 10)

# ---------------------------------------------------------------------------
# Slide 6 — KubeCon EU 2026 Deep Dive
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header_bar(s, "KubeCon + CloudNativeCon Europe 2026",
           "Amsterdam 🇳🇱 · 23–26 March 2026 · CNCF flagship")

add_text(s, Inches(0.5), Inches(1.3), Inches(8.5), Inches(0.6),
         "Speaker session", size=18, bold=True, color=NAVY)
add_text(
    s, Inches(0.5), Inches(1.85), Inches(8.5), Inches(0.9),
    "Lessons Learned Orchestrating Multi-Tenant GPUs on OpenShift AI\n"
    "with NVIDIA KAI on G/H200  —  Luca Berton, Dell Technologies",
    size=14, bold=True, color=ACCENT,
)
add_bullets(
    s, Inches(0.5), Inches(2.85), Inches(8.5), Inches(3.0),
    [
        "Tenant isolation patterns & scheduling on heterogeneous nodes",
        "MIG vs full-GPU trade-offs, throughput vs latency tuning",
        "Driver / firmware pitfalls and safe upgrade & rollback strategies",
        "Day-2 ops: observability, autoscaling and chargeback",
    ],
    size=12,
)

add_text(s, Inches(0.5), Inches(5.6), Inches(8.5), Inches(0.4),
         "Coverage produced around the event", size=14, bold=True, color=NAVY)
add_bullets(
    s, Inches(0.5), Inches(6.0), Inches(8.5), Inches(1.2),
    [
        "Series of speaker & vendor interviews on lucaberton.com/blog",
        "MC role at Cloud Native Rejekts EU 2026 (community unconference)",
        "Side-events guide and RSVP hub published pre-event",
    ],
    size=11,
)

# Right-side info card
add_rect(s, Inches(9.4), Inches(1.3), Inches(3.4), Inches(5.7),
         NAVY)
add_text(s, Inches(9.6), Inches(1.45), Inches(3.0), Inches(0.5),
         "Why it matters", size=14, bold=True, color=WHITE)
add_text(
    s, Inches(9.6), Inches(1.95), Inches(3.0), Inches(5.0),
    "KubeCon EU is the largest cloud-native\n"
    "gathering in Europe with 12,000+\n"
    "attendees from the Kubernetes,\n"
    "AI/ML and platform engineering\n"
    "communities.\n\n"
    "Luca brings a practitioner's lens to\n"
    "GPU orchestration — a top-of-mind\n"
    "topic for every enterprise rolling\n"
    "out generative-AI platforms.",
    size=11, color=WHITE,
)
footer(s, 6, 10)

# ---------------------------------------------------------------------------
# Slide 7 — FOSDEM 2026 Deep Dive
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header_bar(s, "FOSDEM 2026",
           "Brussels 🇧🇪 · 31 Jan – 1 Feb 2026 · Europe's biggest open-source event")

add_text(s, Inches(0.5), Inches(1.3), Inches(8.5), Inches(0.5),
         "On the ground", size=18, bold=True, color=NAVY)
add_bullets(
    s, Inches(0.5), Inches(1.85), Inches(8.5), Inches(2.6),
    [
        "8,000+ attendees · 1,013 talks across 65 dev rooms",
        "Met DeepComputing & RISC-V International Foundation",
        "Themes covered: open silicon, edge AI, hardware/software co-design",
        "Connected with maintainers across CNCF, Ansible & Linux ecosystems",
    ],
    size=12,
)

add_text(s, Inches(0.5), Inches(4.4), Inches(8.5), Inches(0.5),
         "Content produced", size=18, bold=True, color=NAVY)
add_bullets(
    s, Inches(0.5), Inches(4.95), Inches(8.5), Inches(2.0),
    [
        "Blog recap: lucaberton.com/blog/fosdem-2026/",
        "LinkedIn posts with photos & community highlights",
        "Continuation of the FOSDEM 2025 vlog series (BGE on the Road)",
    ],
    size=12,
)

add_rect(s, Inches(9.4), Inches(1.3), Inches(3.4), Inches(5.7), LIGHT)
add_text(s, Inches(9.6), Inches(1.45), Inches(3.0), Inches(0.5),
         "Why FOSDEM", size=14, bold=True, color=NAVY)
add_text(
    s, Inches(9.6), Inches(1.95), Inches(3.0), Inches(5.0),
    "FOSDEM is the unmissable\n"
    "European meeting point for the\n"
    "open-source community.\n\n"
    "Luca's coverage gives sponsors\n"
    "and projects an authentic,\n"
    "practitioner-led narrative —\n"
    "from the dev rooms to the\n"
    "famous beer event.",
    size=11, color=DARK,
)
footer(s, 7, 10)

# ---------------------------------------------------------------------------
# Slide 8 — DevWorld Deep Dive
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header_bar(s, "DevWorld Conference",
           "Amsterdam 🇳🇱 · February 2025 · Europe's largest developer festival")

add_text(s, Inches(0.5), Inches(1.3), Inches(8.5), Inches(0.5),
         "Participation", size=18, bold=True, color=NAVY)
add_bullets(
    s, Inches(0.5), Inches(1.85), Inches(8.5), Inches(2.6),
    [
        "Attendee & networker across the AI, cloud and DevEx tracks",
        "Met framework maintainers, DevRel teams and startup founders",
        "Captured emerging trends in agentic AI and developer tooling",
    ],
    size=12,
)

add_text(s, Inches(0.5), Inches(4.4), Inches(8.5), Inches(0.5),
         "Coverage produced", size=18, bold=True, color=NAVY)
add_bullets(
    s, Inches(0.5), Inches(4.95), Inches(8.5), Inches(2.0),
    [
        "LinkedIn recap with key takeaways for engineering leaders",
        "Blog notes feeding the AI / Kubernetes content series",
        "Cross-promotion in the Ansible Pilot newsletter & site",
    ],
    size=12,
)

add_rect(s, Inches(9.4), Inches(1.3), Inches(3.4), Inches(5.7), NAVY)
add_text(s, Inches(9.6), Inches(1.45), Inches(3.0), Inches(0.5),
         "Why DevWorld", size=14, bold=True, color=WHITE)
add_text(
    s, Inches(9.6), Inches(1.95), Inches(3.0), Inches(5.0),
    "DevWorld brings together\n"
    "developers, DevRel teams and\n"
    "tech founders from across\n"
    "Europe.\n\n"
    "It is where new dev tools and\n"
    "AI products are tested in front\n"
    "of an opinionated audience —\n"
    "and Luca helps amplify the\n"
    "ones that matter.",
    size=11, color=WHITE,
)
footer(s, 8, 10)

# ---------------------------------------------------------------------------
# Slide 9 — Other 2025–2026 Highlights
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header_bar(s, "Other 2025–2026 Highlights",
           "Year-round presence across the cloud-native, AI and open-source calendar")

highlights = [
    ("KubeCon EU 2025 — London 🇬🇧",
     "Apr 2025 · CNCF flagship attendee & community contributor"),
    ("Red Hat Summit: Connect — Utrecht 🇳🇱",
     "Oct 2025 · Sessions on unified AI & application platform"),
    ("Dutch Cloud Native Day — Utrecht 🇳🇱",
     "Jul 2025 · Volunteer & book signing for Kubernetes Recipes"),
    ("FikaWorks Day — Amsterdam 🇳🇱",
     "Oct 2025 · Talk: How to Write a Technical Book and Sell Worldwide"),
    ("CfgMgmtCamp — Ghent 🇧🇪",
     "Feb 2025 · Talk on Ansible + Neo4j GenAI automation"),
    ("FOSDEM 2025 — Brussels 🇧🇪",
     "Feb 2025 · BGE on the Road vlog series"),
    ("BASE Conference — Amsterdam 🇳🇱",
     "Oct 2024 · Panel on AI for SMEs aligned with values"),
    ("TNW Conference — Amsterdam 🇳🇱",
     "Jun 2025 · Tech festival networking"),
]
x0 = Inches(0.5)
y0 = Inches(1.35)
cw = Inches(6.15)
ch = Inches(1.25)
gap = Inches(0.15)
for i, (title, sub) in enumerate(highlights):
    col = i % 2
    row = i // 2
    x = x0 + col * (cw + gap)
    y = y0 + row * (ch + Inches(0.1))
    add_rect(s, x, y, cw, ch, WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, x, y, Inches(0.12), ch, ACCENT)
    add_text(s, x + Inches(0.3), y + Inches(0.18), cw - Inches(0.4),
             Inches(0.5), title, size=13, bold=True, color=NAVY)
    add_text(s, x + Inches(0.3), y + Inches(0.65), cw - Inches(0.4),
             Inches(0.55), sub, size=11, color=GREY)
footer(s, 9, 10)

# ---------------------------------------------------------------------------
# Slide 10 — Work Together / Contact
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(6.7), SW, Inches(0.12), ACCENT)

if LOGO.exists():
    s.shapes.add_picture(str(LOGO), Inches(0.55), Inches(0.5),
                         height=Inches(0.85))

add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.9),
         "Let's work together", size=46, bold=True, color=WHITE)
add_text(
    s, Inches(0.6), Inches(2.85), Inches(12), Inches(0.7),
    "Speaking · Sponsored coverage · Advisory · Workshops",
    size=20, color=ACCENT2,
)

add_rect(s, Inches(0.6), Inches(3.8), Inches(6.0), Inches(2.6),
         RGBColor(0x14, 0x2B, 0x4F))
add_text(s, Inches(0.85), Inches(3.95), Inches(5.5), Inches(0.5),
         "Engagement formats", size=14, bold=True, color=WHITE)
add_bullets(
    s, Inches(0.85), Inches(4.45), Inches(5.5), Inches(2.0),
    [
        "Keynotes, panels & hands-on workshops",
        "Sponsored event coverage (blog, video, social)",
        "Pre-event interviews and recap content",
        "Advisory on AI, GPU & platform engineering",
    ],
    size=12, color=WHITE,
)

add_rect(s, Inches(6.8), Inches(3.8), Inches(6.0), Inches(2.6),
         RGBColor(0x14, 0x2B, 0x4F))
add_text(s, Inches(7.05), Inches(3.95), Inches(5.5), Inches(0.5),
         "Get in touch", size=14, bold=True, color=WHITE)
add_bullets(
    s, Inches(7.05), Inches(4.45), Inches(5.5), Inches(2.0),
    [
        "luca@lucaberton.com",
        "+31 6 1225 5399  ·  Amsterdam, NL",
        "lucaberton.com  ·  ansiblepilot.com",
        "linkedin.com/in/lucaberton  ·  @TheLucaBerton",
    ],
    size=12, color=WHITE,
)

add_text(
    s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
    "Please share: audience size · event date · format (in-person / virtual) · "
    "topic(s) · desired outcomes.",
    size=11, color=RGBColor(0xC9, 0xD7, 0xF0),
)

# Save
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"Wrote {OUT}")
