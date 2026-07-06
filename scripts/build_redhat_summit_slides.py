"""Build the Red Hat Summit 2026 deck by CLONING the KubeCon 2026 slides
(all shapes, rounded rectangles, code blocks, images) onto the Red Hat
Summit light template, so they inherit the Red Hat masters/theme but keep
the original visual layout.
"""
from copy import deepcopy
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
SRC = ROOT / "Media Kit" / "Lessons Learned Orchestrating Multi-Tenant GPUs on OpenShift AI with NVIDIA KAI H200 - Luca Berton Dell Technologies.pptx"
TPL = ROOT / "Media Kit" / "slide" / "Red Hat summit" / "Red Hat Summit 2026 presentation template.pptx"
OUT = ROOT / "Media Kit" / "slide" / "Red Hat summit" / "Lessons Learned Orchestrating Multi-Tenant GPUs on OpenShift AI with NVIDIA KAI H200 - Luca Berton Red Hat Summit 2026.pptx"

TITLE_LAYOUT = "TITLE"
CONTENT_LAYOUT = "CUSTOM_4_17"


def remove_all_slides(prs):
    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    rid_attr = qn("r:id")
    rids = [sldId.get(rid_attr) for sldId in list(sldIdLst)]
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)
    for rid in rids:
        try:
            prs.part.drop_rel(rid)
        except KeyError:
            pass


def find_layout(prs, name):
    for l in prs.slide_layouts:
        if l.name == name:
            return l
    raise KeyError(name)


def clear_placeholders(slide):
    spTree = slide.shapes._spTree
    for sp in list(spTree.iterchildren()):
        if sp.find('.//' + qn('p:ph')) is not None:
            spTree.remove(sp)


def _resolve_placeholder_xfrm(src_layout, ph_idx, ph_type):
    """Walk the source layout (and master) to find an explicit xfrm for the
    placeholder with the given idx/type — so cloned shapes carry their
    original position instead of inheriting it from the destination layout."""
    for spTree in (src_layout.shapes._spTree, src_layout.slide_master.shapes._spTree):
        for sp in spTree.iter(qn('p:sp')):
            ph = sp.find('.//' + qn('p:ph'))
            if ph is None:
                continue
            idx = ph.get('idx', '0')
            typ = ph.get('type', 'body')
            if (ph_idx is not None and idx == str(ph_idx)) or (ph_type and typ == ph_type):
                xfrm = sp.find('.//' + qn('p:spPr') + '/' + qn('a:xfrm'))
                if xfrm is not None:
                    return deepcopy(xfrm)
    return None


def _bake_placeholder_position(new_el, src_slide):
    """If new_el is a placeholder shape lacking explicit xfrm, inject one from
    the source layout/master, and remove the <p:ph> tag so it becomes a free
    shape positioned absolutely on the destination slide."""
    if new_el.tag != qn('p:sp'):
        return
    ph = new_el.find('.//' + qn('p:ph'))
    if ph is None:
        return
    ph_type = ph.get('type')
    spPr = new_el.find(qn('p:spPr'))
    has_xfrm = spPr is not None and spPr.find(qn('a:xfrm')) is not None
    if not has_xfrm:
        xfrm = _resolve_placeholder_xfrm(
            src_slide.slide_layout,
            ph.get('idx'),
            ph_type,
        )
        if xfrm is not None and spPr is not None:
            spPr.insert(0, xfrm)
    # If this is a TITLE placeholder, force the text colour to dark — the
    # source deck used bg1 (white) for titles on a dark master, which would
    # be invisible on the Red Hat light template.
    if ph_type == 'title':
        for rPr in new_el.iter(qn('a:rPr')):
            for sf in rPr.findall(qn('a:solidFill')):
                rPr.remove(sf)
        for defRPr in new_el.iter(qn('a:defRPr')):
            for sf in defRPr.findall(qn('a:solidFill')):
                defRPr.remove(sf)
    # Drop <p:ph> so the cloned shape doesn't try to inherit from dst layout
    ph_parent = ph.getparent()
    ph_parent.remove(ph)


def clone_shapes(src_slide, dst_slide):
    src_part = src_slide.part
    dst_part = dst_slide.part
    dst_spTree = dst_slide.shapes._spTree
    src_spTree = src_slide.shapes._spTree

    # Only copy media / links — NOT slideLayout / slideMaster / theme, so the
    # cloned shapes inherit the Red Hat template look instead of the KubeCon
    # dark master.
    IMAGE_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
    OTHER_ALLOWED = ("hyperlink", "chart", "diagram", "oleObject", "media", "audio", "video")
    rid_map = {}
    for rid, rel in src_part.rels.items():
        if rel.reltype == IMAGE_RELTYPE:
            # Re-import image with a fresh partname inside the destination
            # package to avoid duplicate-name collisions in the output zip.
            blob = rel.target_part.blob
            image_part, new_rid = dst_part.get_or_add_image_part(BytesIO(blob))
            rid_map[rid] = new_rid
        elif any(tok in rel.reltype.lower() for tok in OTHER_ALLOWED):
            if rel.is_external:
                new_rid = dst_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
            else:
                new_rid = dst_part.relate_to(rel.target_part, rel.reltype)
            rid_map[rid] = new_rid

    skip_tags = {qn('p:nvGrpSpPr'), qn('p:grpSpPr'), qn('p:extLst')}
    for child in src_spTree.iterchildren():
        if child.tag in skip_tags:
            continue
        new_el = deepcopy(child)
        for attr in ('r:embed', 'r:link'):
            for el in new_el.iter():
                v = el.get(qn(attr))
                if v and v in rid_map:
                    el.set(qn(attr), rid_map[v])
        for el in new_el.iter(qn('a:hlinkClick')):
            v = el.get(qn('r:id'))
            if v and v in rid_map:
                el.set(qn('r:id'), rid_map[v])
        _bake_placeholder_position(new_el, src_slide)
        dst_spTree.append(new_el)


def main():
    src_prs = Presentation(str(SRC))
    prs = Presentation(str(TPL))
    remove_all_slides(prs)

    title_layout = find_layout(prs, TITLE_LAYOUT)
    content_layout = find_layout(prs, CONTENT_LAYOUT)

    for i, src_slide in enumerate(src_prs.slides):
        layout = title_layout if i == 0 else content_layout
        new_slide = prs.slides.add_slide(layout)
        clear_placeholders(new_slide)
        clone_shapes(src_slide, new_slide)

    # Patch venue mention on first slide
    first = prs.slides[0]
    for sh in first.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if "KubeCon Europe 2026" in r.text:
                    r.text = r.text.replace("KubeCon Europe 2026", "Red Hat Summit 2026")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved {len(prs.slides)} slides → {OUT}")


if __name__ == "__main__":
    main()
