"""Copy speaker notes from the KubeCon source deck into the (already edited)
Red Hat Summit deck, slide-by-slide, in place.

Safe to re-run: it overwrites notes on slides that have a matching index in
the source. It does NOT touch any slide shapes or layout.
"""
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
SRC = ROOT / "Media Kit" / "Lessons Learned Orchestrating Multi-Tenant GPUs on OpenShift AI with NVIDIA KAI H200 - Luca Berton Dell Technologies.pptx"
DST = ROOT / "Media Kit" / "slide" / "Red Hat summit" / "Lessons Learned Orchestrating Multi-Tenant GPUs on OpenShift AI with NVIDIA KAI H200 - Luca Berton Red Hat Summit 2026.pptx"


def get_notes_text_xml(slide):
    """Return the <p:txBody> element of the notes placeholder, or None."""
    if not slide.has_notes_slide:
        return None
    notes = slide.notes_slide
    for sp in notes.shapes._spTree.iter(qn('p:sp')):
        ph = sp.find('.//' + qn('p:ph'))
        if ph is not None and ph.get('type') == 'body':
            return sp.find(qn('p:txBody'))
    # Fallback: first sp with txBody
    for sp in notes.shapes._spTree.iter(qn('p:sp')):
        tb = sp.find(qn('p:txBody'))
        if tb is not None:
            return tb
    return None


def main():
    src = Presentation(str(SRC))
    dst = Presentation(str(DST))

    n = min(len(src.slides), len(dst.slides))
    copied = 0
    for i in range(n):
        src_tb = get_notes_text_xml(src.slides[i])
        if src_tb is None:
            continue
        # Ensure destination notes slide exists
        dst_notes = dst.slides[i].notes_slide  # creates if missing
        dst_tb = get_notes_text_xml(dst.slides[i])
        if dst_tb is None:
            continue
        parent = dst_tb.getparent()
        new_tb = deepcopy(src_tb)
        parent.replace(dst_tb, new_tb)
        copied += 1
        # Preview
        text_lines = []
        for p in new_tb.iter(qn('a:p')):
            txt = "".join(t.text or "" for t in p.iter(qn('a:t')))
            if txt.strip():
                text_lines.append(txt)
        snippet = " | ".join(text_lines)[:90]
        print(f"  slide {i+1}: {snippet}")

    dst.save(str(DST))
    print(f"\nCopied notes onto {copied}/{n} slides → {DST.name}")


if __name__ == "__main__":
    main()
